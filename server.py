"""Tradewind 本地服务：FastAPI 提供 API + 托管 React 前端，单进程单端口。

exe 化基础：打包后由 run.py 启动本服务，浏览器访问 http://127.0.0.1:8101
"""

from __future__ import annotations

import csv
import base64
import hashlib
import io
import json
import os
import platform
import re
import shutil
import sys
import tempfile
import threading
import time
import uuid
import warnings
import zipfile
from collections import Counter
from pathlib import Path
from urllib.parse import urlsplit
from xml.etree import ElementTree

# markitdown 的音频模块（pydub）无 ffmpeg 时打警告；与文档/OCR 功能无关，静默掉
warnings.filterwarnings("ignore", message="Couldn't find ffmpeg or avconv.*")

import httpx
from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from app.config import (
    PROVIDER_PRESETS, ROOT_DIR, VISION_PROVIDERS, get_active_provider, get_company_profile,
    get_provider_config, get_vision_config, get_vision_failover_configs,
    get_vision_provider_config, set_active_provider, set_api_key,
    set_company_profile, set_provider_params, set_vision_config, settings, vision_key_format_error,
)
from app.email_agent import GenerationCancelled, generate_email
from app.crawler.progress import use_progress_sink
from app.crawler.result_utils import matches_targets
from app.memory import delete_email, delete_emails, recent_emails
from app.task_utils import prune_finished_tasks

ROOT = ROOT_DIR
APP_VERSION = "0.2.0"


def _resolve_frontend_dist() -> Path | None:
    """静态资源目录解析：
    ① PyInstaller 打包：资源解包在 _MEIPASS/frontend/dist（_internal 里）
    ② exe 旁手工放置 frontend/dist（兼容）
    ③ 开发模式：项目根 frontend/dist
    """
    if getattr(sys, "frozen", False):
        meipass = getattr(sys, "_MEIPASS", None)
        if meipass:
            p = Path(meipass) / "frontend" / "dist"
            if p.exists():
                return p
        p = Path(sys.executable).resolve().parent / "frontend" / "dist"
        if p.exists():
            return p
    # ROOT may point to the per-user data directory when run.py isolates
    # mutable desktop data. Source assets still live beside server.py.
    p = Path(__file__).resolve().parent / "frontend" / "dist"
    return p if p.exists() else None


FRONTEND_DIST = _resolve_frontend_dist()
DATA_DIR = ROOT / "data"
PRODUCTS_FILE = DATA_DIR / "products.json"
TEMPLATES_FILE = DATA_DIR / "emails.json"
CUSTOMERS_FILE = DATA_DIR / "customers.json"
MEMORY_DB_FILE = DATA_DIR / "tradewind_memory.db"
CRAWLER_SEEN_FILE = DATA_DIR / "crawler_seen.json"
CRAWLER_PHOTOS_DIR = DATA_DIR / "crawler_photos"
PHOTO_SCAN_DIR = DATA_DIR / "photo_scan"


def _ensure_data_files() -> None:
    """exe 首启数据模板复制：PyInstaller 把 data 解包在 _MEIPASS/data，
    而运行时的 ROOT/data 在 exe 旁。目标文件不存在时从解包目录复制过去，
    保证业务用户电脑首启就有产品/话术模板（用户已改的数据不会被覆盖）。"""
    meipass = getattr(sys, "_MEIPASS", None)
    if not meipass:
        return
    src = Path(meipass) / "data"
    if not src.exists():
        return
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    for name in ("products.json", "emails.json"):
        s = src / name
        d = DATA_DIR / name
        if s.exists() and not d.exists():
            try:
                shutil.copy2(s, d)
                print(f"[tradewind] 已复制初始{name}到数据目录 {d}")
            except OSError:
                pass


_ensure_data_files()

app = FastAPI(title="Tradewind", version=APP_VERSION)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


# ---------- 工具 ----------

def _read_json(path: Path, default: list) -> list:
    if not path.exists():
        return default
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        return data if isinstance(data, list) else default
    except Exception:  # noqa: BLE001
        return default


def _write_json(path: Path, data: list) -> None:
    """原子写入 JSON，避免进程中断时留下半截用户数据。"""
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_name(f".{path.name}.{uuid.uuid4().hex}.tmp")
    try:
        with tmp.open("w", encoding="utf-8", newline="") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
            f.flush()
            os.fsync(f.fileno())
        os.replace(tmp, path)
    finally:
        tmp.unlink(missing_ok=True)


def _next_id(items: list, prefix: str) -> str:
    n = 1
    for it in items:
        m = str(it.get("id", ""))
        if m.startswith(prefix) and m[len(prefix):].isdigit():
            n = max(n, int(m[len(prefix):]) + 1)
    return f"{prefix}{n:03d}"


def _find_template(template_id: str) -> dict | None:
    """按当前运行数据目录查找模板，单封与批量共用同一条链路。"""
    if not template_id:
        return None
    return next(
        (item for item in _read_json(TEMPLATES_FILE, [])
         if str(item.get("id")) == template_id),
        None,
    )


def _parse_product_doc(text: str) -> list[dict]:
    """从文档转换出的 Markdown/文本里启发式切分产品条目：
    按空行分块，块首行作标题，其余作描述。"""
    items: list[dict] = []
    for block in re.split(r"\n\s*\n", text):
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        title = lines[0].strip(" #*-=~").rstrip("：:|")
        if len(title) < 2 or len(title) > 60:
            continue
        content = "\n".join(lines[1:]).strip()
        if not content:
            content = "（从文档提取，无详细描述）"
        items.append({"title": title[:60], "content": content[:2000]})
        if len(items) >= 30:
            break
    return items


# ---------- 模型 ----------

class EmailRequest(BaseModel):
    customer: str
    country: str = ""
    product: str = "医美设备"
    extra: str = ""
    judge: bool = False
    language: str = "zh-hant"   # zh-hant（香港繁体）| en
    format: str = "email"       # email | whatsapp
    template_id: str = ""


class EmailResponse(BaseModel):
    customer: str
    country: str
    product: str
    email: str
    issues: list[str] = []
    scores: dict | None = None
    tokens: int = 0
    time_s: float = 0.0
    templates_used: list[str] = []


class ConfigRequest(BaseModel):
    provider: str = ""
    api_key: str = ""


class ActivateRequest(BaseModel):
    provider: str = ""


class ProviderParamsRequest(BaseModel):
    provider: str = ""
    base_url: str = ""
    model: str = ""


class CompanyProfileRequest(BaseModel):
    sender_name: str = ""
    company_name: str = ""
    email: str = ""
    whatsapp: str = ""
    website: str = ""


class ProductRequest(BaseModel):
    title: str
    content: str
    tags: list[str] = []


class CustomerRequest(BaseModel):
    name: str
    country: str = ""
    city: str = ""
    email: str = ""
    website: str = ""
    phone: str = ""
    notes: str = ""
    status: str = "new"


class BatchRequest(BaseModel):
    rows: list[dict]
    product: str = "医美设备"
    judge: bool = False
    language: str = "zh-hant"
    format: str = "email"
    template_id: str = ""


# ---------- 配置 API ----------

@app.get("/api/config")
def get_config() -> dict:
    """前端探测：当前是否已配置 key + 各 provider 状态（不回显 key 本身）。"""
    active = get_active_provider()
    providers = []
    for pid, preset in PROVIDER_PRESETS.items():
        cfg = get_provider_config(pid)
        providers.append({
            "id": pid,
            "name": preset["name"],
            "base_url": cfg["base_url"],
            "model": cfg["model"],
            "has_key": bool(cfg["api_key"]),
        })
    acfg = get_provider_config(active)
    vision_current = get_vision_config()
    vision_candidates = get_vision_failover_configs(vision_current["provider"])
    vision_effective = vision_candidates[0] if vision_candidates else vision_current
    return {
        "has_key": bool(acfg["api_key"]),
        "model": acfg["model"],
        "active_provider": active,
        "providers": providers,
        "company_profile": get_company_profile(),
        # 视觉识别（爬虫照片分析；不回显 key，只报状态）
        "vision": {
            "configured": bool(vision_candidates),
            "provider": vision_current["provider"],
            "model": vision_current["model"],
            "effective_provider": vision_effective["provider"] if vision_candidates else "",
            "effective_model": vision_effective["model"] if vision_candidates else "",
            "providers": {
                pid: {
                    "name": pv["name"],
                    "models": pv["models"],
                    "has_key": bool(get_vision_provider_config(pid)["api_key"]),
                    "model": get_vision_provider_config(pid)["model"],
                }
                for pid, pv in VISION_PROVIDERS.items()
            },
        },
        # 数据存放位置（桌面版 = AppLocalData；开发版 = 项目 data）
        "data_dir": str(DATA_DIR),
        "data_files": {
            "customers": str(CUSTOMERS_FILE),
            "products": str(PRODUCTS_FILE),
            "templates": str(TEMPLATES_FILE),
            "memory": str(MEMORY_DB_FILE),
            "crawler_csv": str(DATA_DIR / "maps_hk.csv"),
        },
    }


@app.post("/api/config")
def post_config(req: ConfigRequest) -> dict:
    """保存某 provider 的 key 并验证（401 = key 无效；网络错误不阻塞保存），保存后设为激活。"""
    pid = req.provider or get_active_provider()
    if pid not in PROVIDER_PRESETS:
        return {"ok": False, "error": "未知服务商"}
    key = req.api_key.strip()
    if not key:
        return {"ok": False, "error": "API Key 不能为空"}
    base = get_provider_config(pid)["base_url"] or "https://api.deepseek.com"
    verified: bool | None = None
    try:
        resp = httpx.get(
            f"{base}/models",
            headers={"Authorization": f"Bearer {key}"},
            timeout=10,
        )
        if resp.status_code == 401:
            return {"ok": False, "error": "Key 无效（401），请检查后重填"}
        verified = resp.status_code == 200
    except Exception:  # noqa: BLE001 - 网络问题不阻塞保存
        verified = None
    set_api_key(pid, key)
    set_active_provider(pid)
    return {"ok": True, "has_key": True, "verified": verified}


@app.post("/api/config/activate")
def activate_provider(req: ActivateRequest) -> dict:
    """切换当前生效的 provider（不需要重填 key）。"""
    if req.provider not in PROVIDER_PRESETS:
        return {"ok": False, "error": "未知服务商"}
    if not get_provider_config(req.provider)["api_key"]:
        return {"ok": False, "error": "该服务商尚未配置 Key"}
    set_active_provider(req.provider)
    return {"ok": True, "active_provider": req.provider}


@app.post("/api/config/params")
def provider_params(req: ProviderParamsRequest) -> dict:
    """自定义服务商：保存 base_url / model。"""
    if req.provider not in PROVIDER_PRESETS:
        return {"ok": False, "error": "未知服务商"}
    set_provider_params(req.provider, req.base_url, req.model)
    return {"ok": True}


@app.post("/api/config/company")
def company_profile(req: CompanyProfileRequest) -> dict:
    """保存本机发件人及公司签名资料。"""
    if not req.company_name.strip():
        return {"ok": False, "error": "公司名称不能为空"}
    set_company_profile(req.model_dump())
    return {"ok": True, "company_profile": get_company_profile()}


class VisionRequest(BaseModel):
    api_key: str = ""
    model: str = ""
    provider: str = ""  # glm | qwen | volc


@app.post("/api/config/vision")
def post_vision_config(req: VisionRequest) -> dict:
    """保存视觉识别配置（爬虫照片分析专用；写开发信仍走 provider）。"""
    if req.provider and req.provider not in VISION_PROVIDERS:
        return {"ok": False, "error": "未知视觉服务商"}
    provider = req.provider or get_vision_config()["provider"]
    effective_key = req.api_key.strip() or get_vision_provider_config(provider)["api_key"]
    key_error = vision_key_format_error(provider, effective_key)
    if key_error:
        return {"ok": False, "error": key_error}
    set_vision_config(req.provider, req.api_key, req.model)
    v = get_vision_config()
    return {"ok": True, "configured": bool(v["api_key"]), "provider": v["provider"], "model": v["model"]}


_SNAP_RE = re.compile(r"-\d{4}-\d{2}-\d{2}(?:-\d+)?$")
_CHAT_PREFIX = ("gpt-",)
_DROP_KEY = ("realtime", "audio", "tts", "whisper", "dall-e", "embedding", "moderation",
             "fine-tune", "search", "image", "video", "transcription", "translation", "speech")


def _slim_models(models: list[str], provider: str) -> list[str]:
    """OpenAI /models 返回 118 个（含 embedding/音频/快照），前端下拉太爆炸：
    只留常用聊天模型，去快照版本，白名单优先排序。"""
    if provider != "openai":
        return models
    seen: set[str] = set()
    out: list[str] = []
    for m in models:
        low = m.lower()
        base = _SNAP_RE.sub("", low)  # gpt-4.1-mini-2025-07-18 → gpt-4.1-mini
        if not base.startswith(_CHAT_PREFIX) or any(k in low for k in _DROP_KEY):
            continue
        if base not in seen:
            seen.add(base)
            out.append(base)
    order = ["gpt-4.1", "gpt-4.1-mini", "gpt-4.1-nano",
             "gpt-4o", "gpt-4o-mini", "gpt-5", "gpt-5.2"]
    out.sort(key=lambda m: (order.index(m) if m in order else 99, m))
    return out[:10]


@app.get("/api/config/{provider}/models")
def provider_models(provider: str) -> dict:
    """拉取某服务商当前 Key 可用的模型列表（调各家 /models 接口，免费）。"""
    if provider not in PROVIDER_PRESETS:
        return {"ok": False, "error": "未知服务商", "models": []}
    cfg = get_provider_config(provider)
    if not cfg["api_key"]:
        return {"ok": False, "error": "未配置 Key", "models": []}
    base = cfg["base_url"] or "https://api.deepseek.com"
    # OpenAI 接口国内直连不通：模型列表拉取走代理（DeepSeek/Kimi 国内直连无需代理）
    proxy = settings.crawler_proxy if "openai.com" in base else None
    try:
        resp = httpx.get(
            f"{base}/models",
            headers={"Authorization": f"Bearer {cfg['api_key']}"},
            timeout=15,
            proxy=proxy or None,
        )
        if resp.status_code != 200:
            return {"ok": False, "error": f"获取失败（{resp.status_code}）", "models": []}
        data = resp.json()
        models = [m["id"] for m in data.get("data", []) if isinstance(m, dict) and m.get("id")]
        return {"ok": True, "models": _slim_models(models, provider)}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc), "models": []}


# ---------- 业务 API：单封生成 ----------

@app.get("/api/health")
def health() -> dict:
    return {"status": "ok", "service": "tradewind", "version": APP_VERSION}


def _diagnostic_file_stats(path: Path, *, count_json: bool = False) -> dict:
    stats: dict = {"exists": path.is_file(), "size_bytes": 0}
    if not stats["exists"]:
        if count_json:
            stats["item_count"] = 0
        return stats
    try:
        stats["size_bytes"] = path.stat().st_size
    except OSError:
        pass
    if count_json:
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
            stats["item_count"] = len(payload) if isinstance(payload, list) else 0
            stats["readable"] = isinstance(payload, list)
        except Exception:  # noqa: BLE001 - 诊断只报告可读性，不返回原文件内容
            stats["item_count"] = 0
            stats["readable"] = False
    return stats


def _diagnostic_error_summary(path: Path) -> dict:
    """仅统计错误类型/HTTP 状态，绝不返回可能含业务信息的原始日志。"""
    summary: dict = {"exists": path.is_file(), "size_bytes": 0, "error_types": {}, "http_statuses": {}}
    if not summary["exists"]:
        return summary
    try:
        raw = path.read_bytes()
        summary["size_bytes"] = len(raw)
        text = raw[-262_144:].decode("utf-8", errors="replace")
        error_types = Counter(re.findall(r"\b[A-Za-z_][A-Za-z0-9_]*(?:Error|Exception|Timeout)\b", text))
        statuses = Counter(re.findall(r"\b(?:HTTP\s*)?([45]\d{2})\b", text, flags=re.I))
        summary["error_types"] = dict(error_types.most_common(20))
        summary["http_statuses"] = dict(statuses.most_common(20))
    except OSError:
        summary["readable"] = False
    return summary


def _diagnostic_data_mode() -> str:
    if not getattr(sys, "frozen", False):
        return "source"
    try:
        program_dir = Path(sys.executable).resolve().parent
        return "folder-portable" if ROOT.resolve() == program_dir else "appdata"
    except OSError:
        return "packaged"


def _build_diagnostic_report() -> dict:
    active_provider = get_active_provider()
    providers = {
        pid: {
            "configured": bool(get_provider_config(pid).get("api_key")),
            "model": str(get_provider_config(pid).get("model") or ""),
        }
        for pid in PROVIDER_PRESETS
    }
    vision = get_vision_config()
    vision_providers = {
        pid: {
            "configured": bool(get_vision_provider_config(pid).get("api_key")),
            "model": str(get_vision_provider_config(pid).get("model") or ""),
        }
        for pid in VISION_PROVIDERS
    }
    photo_files = []
    if CRAWLER_PHOTOS_DIR.is_dir():
        try:
            photo_files = [
                path for path in CRAWLER_PHOTOS_DIR.rglob("*")
                if path.is_file() and path.suffix.lower() in PHOTO_SUFFIXES
            ]
        except OSError:
            photo_files = []

    return {
        "schema_version": 1,
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        "application": {
            "name": "Tradewind",
            "version": APP_VERSION,
            "packaged": bool(getattr(sys, "frozen", False)),
            "data_mode": _diagnostic_data_mode(),
        },
        "runtime": {
            "os": platform.system(),
            "os_release": platform.release(),
            "architecture": platform.machine(),
            "python": platform.python_version(),
            "data_writable": os.access(DATA_DIR, os.W_OK),
            "frontend_available": FRONTEND_DIST is not None,
        },
        "models": {
            "active_provider": active_provider,
            "providers": providers,
            "vision_provider": vision.get("provider", ""),
            "vision_model": vision.get("model", ""),
            "vision_providers": vision_providers,
        },
        "data": {
            "customers": _diagnostic_file_stats(CUSTOMERS_FILE, count_json=True),
            "products": _diagnostic_file_stats(PRODUCTS_FILE, count_json=True),
            "templates": _diagnostic_file_stats(TEMPLATES_FILE, count_json=True),
            "history_db": _diagnostic_file_stats(MEMORY_DB_FILE),
            "crawler_photos": {
                "store_count": sum(1 for path in CRAWLER_PHOTOS_DIR.iterdir() if path.is_dir()) if CRAWLER_PHOTOS_DIR.is_dir() else 0,
                "photo_count": len(photo_files),
            },
        },
        "errors": {
            "crawler": _diagnostic_error_summary(DATA_DIR / "crawler_errors.log"),
            "document_import": _diagnostic_error_summary(DATA_DIR / "ocr_errors.log"),
        },
        "privacy": {
            "raw_logs_included": False,
            "keys_included": False,
            "customer_content_included": False,
            "photos_included": False,
            "email_content_included": False,
        },
    }


@app.get("/api/diagnostics/export")
def export_diagnostics() -> Response:
    report = _build_diagnostic_report()
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as output:
        output.writestr("diagnostic.json", json.dumps(report, ensure_ascii=False, indent=2))
        output.writestr(
            "README.txt",
            "Tradewind 本地诊断包\n\n"
            "此文件由用户主动导出，没有后台上传。\n"
            "包含：版本、运行环境、模型名称、是否已配置 Key、数据数量及脱敏错误类型统计。\n"
            "不包含：API Key、客户资料、邮件正文、搜索词、网址、照片、文件路径或原始日志。\n",
        )
    filename = f"Tradewind-diagnostics-{time.strftime('%Y%m%d-%H%M%S')}.zip"
    return Response(
        content=archive.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.post("/api/email", response_model=EmailResponse)
def api_email(req: EmailRequest) -> EmailResponse:
    """生成一封开发信（同步；judge 打分可选，会慢几秒）。"""
    template_record = _find_template(req.template_id)
    if req.template_id and template_record is None:
        raise ValueError("所选话术模板不存在，请刷新后重试")
    result = generate_email(
        req.customer, req.country, req.product,
        judge=req.judge, verbose=False, extra=req.extra,
        language=req.language, format_=req.format, template_record=template_record,
    )
    return EmailResponse(**{k: result.get(k) for k in EmailResponse.model_fields})


# ---------- 业务 API：任务化生成（单封/批量共用，支持取消） ----------

TASKS: dict[str, dict] = {}
_task_lock = threading.Lock()


def _new_task(total: int) -> str:
    task_id = uuid.uuid4().hex[:12]
    with _task_lock:
        prune_finished_tasks(TASKS)
        TASKS[task_id] = {
            "status": "running", "cancel": False,
            "total": total, "done": 0, "current": "",
            "results": [], "result": None, "error": "",
            "created_at": time.time(),
        }
    return task_id


def _cancelled(task_id: str) -> bool:
    return bool(TASKS.get(task_id, {}).get("cancel", False))


def _finish(task_id: str, status: str, **extra) -> None:
    with _task_lock:
        task = TASKS.get(task_id)
        if not task:
            return
        task["status"] = status
        task.update(extra)
        task["finished_at"] = time.time()


def _append_stream(task_id: str, chunk: str) -> None:
    """流式内容实时写入任务（前端轮询取到即打字机展示）。"""
    with _task_lock:
        task = TASKS.get(task_id)
        if task:
            task["stream"] = task.get("stream", "") + chunk


# --- 单封 ---

@app.post("/api/email/start")
def api_email_start(req: EmailRequest) -> dict:
    """创建单封生成任务（后台跑，前端轮询；可取消）。"""
    template_record = _find_template(req.template_id)
    if req.template_id and template_record is None:
        return {"task_id": "", "error": "所选话术模板不存在，请刷新后重试"}
    task_id = _new_task(1)

    def execute() -> None:
        try:
            result = generate_email(
                req.customer, req.country, req.product,
                judge=req.judge, verbose=False, extra=req.extra,
                cancel_check=lambda: _cancelled(task_id),
                stream_callback=lambda c: _append_stream(task_id, c),
                language=req.language, format_=req.format,
                template_record=template_record,
            )
            _finish(task_id, "done", result=result, done=1)
        except GenerationCancelled:
            _finish(task_id, "cancelled")
        except Exception as exc:  # noqa: BLE001
            _finish(task_id, "error", error=str(exc))

    threading.Thread(target=execute, daemon=True).start()
    return {"task_id": task_id}


# --- 批量 ---

def _run_batch(task_id: str, rows: list[dict], product: str, judge: bool,
               language: str = "zh-hant", format_: str = "email",
               template_record: dict | None = None) -> None:
    results: list[dict] = []
    total = len(rows)
    for i, row in enumerate(rows):
        if _cancelled(task_id):
            _finish(task_id, "cancelled", results=results, done=i)
            return
        name = (row.get("name") or "").strip()
        if not name:
            continue
        # 新的一封开始，重置流式缓冲
        with _task_lock:
            TASKS[task_id]["stream"] = ""
        try:
            r = generate_email(
                name, (row.get("country") or "").strip(), product,
                judge=judge, verbose=False, extra=(row.get("notes") or "").strip(),
                cancel_check=lambda: _cancelled(task_id),
                stream_callback=lambda c: _append_stream(task_id, c),
                language=language, format_=format_,
                template_record=template_record,
            )
            results.append({
                "name": name, "country": row.get("country", ""),
                "email": r["email"], "issues": r["issues"],
                "scores": r["scores"], "ok": True,
                "judge_status": r.get("judge_status", "disabled"),
                "templates_used": r.get("templates_used", []),
            })
        except GenerationCancelled:
            _finish(task_id, "cancelled", results=results, done=i)
            return
        except Exception as exc:  # noqa: BLE001 - 单个客户失败不中断批量
            results.append({
                "name": name, "country": row.get("country", ""),
                "email": "", "issues": [], "scores": None,
                "ok": False, "error": str(exc),
            })
        with _task_lock:
            TASKS[task_id]["done"] = i + 1
            TASKS[task_id]["current"] = name
    _finish(task_id, "done", results=results, done=total)


@app.post("/api/email/batch")
def api_batch(req: BatchRequest) -> dict:
    """创建批量生成任务（后台跑，前端轮询；可取消，保留已完成部分）。"""
    if not req.rows:
        return {"task_id": "", "error": "名单为空"}
    template_record = _find_template(req.template_id)
    if req.template_id and template_record is None:
        return {"task_id": "", "error": "所选话术模板不存在，请刷新后重试"}
    task_id = _new_task(len(req.rows))
    threading.Thread(
        target=_run_batch,
        args=(task_id, req.rows, req.product, req.judge, req.language, req.format, template_record),
        daemon=True,
    ).start()
    return {"task_id": task_id}


@app.get("/api/email/tasks/{task_id}")
def api_task(task_id: str) -> dict:
    task = TASKS.get(task_id)
    if not task:
        return {"status": "not_found"}
    return task


@app.post("/api/email/tasks/{task_id}/cancel")
def api_task_cancel(task_id: str) -> dict:
    """请求停止：置取消标记，生成线程下个 chunk 即中止（1-2 秒内）。"""
    with _task_lock:
        task = TASKS.get(task_id)
        if not task:
            return {"ok": False, "error": "任务不存在"}
        if task["status"] != "running":
            return {"ok": False, "error": "任务已结束"}
        task["cancel"] = True
    return {"ok": True}


# ---------- 业务 API：产品资料 ----------

@app.get("/api/products")
def api_products() -> list[dict]:
    return _read_json(PRODUCTS_FILE, [])


@app.post("/api/products")
def api_product_add(req: ProductRequest) -> dict:
    items = _read_json(PRODUCTS_FILE, [])
    item = {
        "id": _next_id(items, "prod-"),
        "title": req.title.strip(),
        "content": req.content.strip(),
        "source": "frontend",
        "tags": req.tags or ["产品"],
    }
    items.append(item)
    _write_json(PRODUCTS_FILE, items)
    return {"ok": True, "item": item}


_OCR_ENGINE = None  # 模块级单例：OCR 引擎只初始化一次（模型加载约 5s）


def _get_ocr():
    global _OCR_ENGINE
    if _OCR_ENGINE is None:
        from rapidocr_onnxruntime import RapidOCR

        _OCR_ENGINE = RapidOCR()
    return _OCR_ENGINE


def _ocr_lines(result) -> list[str]:
    """RapidOCR 结果 → 按阅读顺序的行文本（y 行聚合 + x 左排序）。"""
    if not result:
        return []
    items = []
    for box, text, _score in result:
        ys = [p[1] for p in box]
        xs = [p[0] for p in box]
        items.append((min(ys), min(xs), str(text).strip()))
    items.sort(key=lambda t: (round(t[0] / 24), t[1]))  # 行高约 24px，按行聚合
    return [t for _, _, t in items if t]


def _ocr_file(path: Path, suffix: str, cancel_check=None) -> str:
    """图片型文档 OCR：PDF 逐页渲染识别，图片直接识别。返回按页分隔的文本。

    cancel_check：每页检查一次（None 则不检查），返回 True 时抛 _ExtractCancelled。
    ImportError 向上抛（调用方给出安装提示）。
    """
    engine = _get_ocr()
    parts: list[str] = []
    if suffix == ".pdf":
        import fitz  # PyMuPDF
        import cv2
        import numpy as np

        doc = fitz.open(str(path))
        try:
            for page in doc:
                if cancel_check and cancel_check():
                    raise _ExtractCancelled()
                # 内存管道：pix → PNG bytes → ndarray → OCR（不落临时文件，
                # 避免 fitz 删除被占用文件在 Windows 上报 Permission denied）
                pix = page.get_pixmap(dpi=150)
                png = pix.tobytes("png")
                mat = cv2.imdecode(np.frombuffer(png, dtype=np.uint8), cv2.IMREAD_COLOR)
                if mat is None:
                    continue
                result, _ = engine(mat)
                lines = _ocr_lines(result)
                if lines:
                    parts.append("\n".join(lines))
        finally:
            doc.close()
    elif suffix in (".png", ".jpg", ".jpeg", ".bmp", ".webp"):
        result, _ = engine(str(path))
        lines = _ocr_lines(result)
        if lines:
            parts.append("\n".join(lines))
    return "\n\n".join(parts)


def _read_text_document(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-16", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_docx_text(path: Path) -> str:
    namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    with zipfile.ZipFile(path) as archive:
        root = ElementTree.fromstring(archive.read("word/document.xml"))
    paragraphs: list[str] = []
    for paragraph in root.iter(f"{namespace}p"):
        line = "".join(node.text or "" for node in paragraph.iter(f"{namespace}t")).strip()
        if line:
            paragraphs.append(line)
    return "\n".join(paragraphs)


def _extract_native_document_text(path: Path, suffix: str) -> str:
    """使用稳定的本地解析器提取常见格式，避免可选插件缺失拖垮全部格式。"""
    if suffix in {".txt", ".md", ".csv"}:
        return _read_text_document(path)
    if suffix in {".html", ".htm"}:
        from bs4 import BeautifulSoup

        return BeautifulSoup(_read_text_document(path), "html.parser").get_text("\n")
    if suffix == ".docx":
        return _extract_docx_text(path)
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            parts: list[str] = []
            for sheet in workbook.worksheets:
                rows = [
                    "\t".join(str(value).strip() for value in row if value is not None)
                    for row in sheet.iter_rows(values_only=True)
                ]
                text = "\n".join(row for row in rows if row)
                if text:
                    parts.append(f"[{sheet.title}]\n{text}")
            return "\n\n".join(parts)
        finally:
            workbook.close()
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(path)
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines = [str(shape.text).strip() for shape in slide.shapes if hasattr(shape, "text") and str(shape.text).strip()]
            if lines:
                slides.append(f"[第 {index} 页]\n" + "\n".join(lines))
        return "\n\n".join(slides)
    if suffix == ".pdf":
        import fitz

        document = fitz.open(str(path))
        try:
            return "\n\n".join(page.get_text("text").strip() for page in document if page.get_text("text").strip())
        finally:
            document.close()
    return ""


def _doc_to_text(tmp: Path, suffix: str, cancel_check=None) -> tuple[str, bool]:
    """文档 → 文本：原生解析优先，MarkItDown 兜底，无文字再 OCR。"""
    text = ""
    try:
        text = _extract_native_document_text(tmp, suffix)
    except Exception:  # noqa: BLE001 - 原生解析失败继续走兼容解析/OCR
        text = ""
    if text.strip():
        return text.strip(), False

    try:
        from markitdown import MarkItDown

        md = MarkItDown()
        res = md.convert(str(tmp))
        text = res.text_content or ""
    except ImportError:
        text = ""  # 无 MarkItDown 时直接走 OCR
    except Exception:  # noqa: BLE001 - 转换失败也走 OCR 兜底
        text = ""
    ocr_used = False
    if not text.strip():
        text = _ocr_file(tmp, suffix, cancel_check=cancel_check)
        ocr_used = True
    return text, ocr_used


_SUMMARY_PROMPTS = {
    "product": (
        "你是医美设备行业的产品资料整理助手。下面是从产品宣传资料提取的文本"
        "（可能含 OCR 噪声、页码、宣传标语）。请总结成一条结构化产品资料。\n\n"
        "要求：\n"
        "1. 标题：品牌+型号+设备类别（如：皮肤检测治疗一体机 DemoMed Vision-100）\n"
        "2. 内容：200-400 字中文（专业术语可保留英文），包含：核心卖点、关键参数"
        "（波长/能量/倍率/模式等数字信息）、适应症或应用场景、认证（如有）、适合的客户类型。\n"
        "3. 标签 tags：3-5 个检索关键词（中文，第一个固定为「产品」，其余覆盖设备类别/技术/适应症，如[\"产品\", \"激光脱毛\", \"808nm\"]）\n"
        "4. 忽略与产品无关的噪声（页眉页脚、装饰性口号）。\n"
        "5. 只输出 JSON：{\"title\": \"...\", \"content\": \"...\", \"tags\": [\"...\"]}\n\n"
    ),
    "template": (
        "你是文档文字提取助手。下面是从历史邮件/聊天记录提取的文本。\n\n"
        "要求：\n"
        "1. content 逐字保留原文，禁止改写、整理、润色、总结或翻译。\n"
        "2. 仅剔除明显噪声（页眉页脚、导航菜单、乱码、装饰符号）。\n"
        "3. title 只概括使用场景，不得写入 content。\n"
        "4. 只输出 JSON：{\"title\": \"...\", \"content\": \"完整原文\", \"tags\": [\"邮件\", \"模板\"]}\n\n"
    ),
}


def _summarize_doc(text: str, fallback_title: str, kind: str = "product", cancel_check=None) -> dict | None:
    """LLM 把识别文本总结成一条结构化条目（产品资料 or 话术模板）。

    cancel_check：调用前检查一次（invoke 无法中途打断，但可在进入前停）。
    失败（未配置 key/网络/输出异常）返回 None，调用方回退原始识别块。
    """
    try:
        if cancel_check and cancel_check():
            raise _ExtractCancelled()
        from app.llm import build_llm

        llm = build_llm(temperature=0.2)
        prompt = _SUMMARY_PROMPTS.get(kind, _SUMMARY_PROMPTS["product"]) + f"识别文本：\n{text[:6000]}"
        resp = llm.invoke(prompt)
        content = resp.content if hasattr(resp, "content") else str(resp)
        m = re.search(r"\{.*\}", content, re.S)
        if m:
            data = json.loads(m.group(0))
            title = str(data.get("title") or "").strip() or fallback_title
            body = str(data.get("content") or "").strip()
            tags = data.get("tags")
            if isinstance(tags, list):
                tags = [str(t).strip() for t in tags if str(t).strip()][:6]
            else:
                tags = []
            if body:
                return {"title": title[:80], "content": body, "tags": tags}
        # 没按 JSON 输出 → 整段作为内容
        clean = content.strip().strip("`").strip()
        if len(clean) >= 20:
            return {"title": fallback_title, "content": clean[:1500], "tags": []}
    except _ExtractCancelled:
        raise
    except Exception:  # noqa: BLE001 - 总结失败回退原始块
        return None
    return None


# ---------- 业务 API：文档识别导入（任务制，支持取消；products/templates 共用） ----------

class _ExtractCancelled(Exception):
    """识别被用户取消（停止按钮触发）。"""


EXTRACT_TASKS: dict[str, dict] = {}


async def _start_extract(file: UploadFile, kind: str) -> dict:
    """接收上传 → 落临时文件 → 后台线程识别。返回 task_id（前端轮询）。"""
    suffix = Path(file.filename or "doc").suffix.lower()
    tmp = Path(tempfile.gettempdir()) / f"tw_{uuid.uuid4().hex}{suffix}"
    tmp.write_bytes(await file.read())
    task_id = uuid.uuid4().hex[:12]
    prune_finished_tasks(EXTRACT_TASKS)
    EXTRACT_TASKS[task_id] = {
        "status": "running", "cancel": False, "error": "",
        "items": [], "mode": "", "ocr": False,
        "filename": file.filename or "", "created_at": time.time(),
    }
    threading.Thread(target=_run_extract, args=(task_id, tmp, kind), daemon=True).start()
    return {"ok": True, "task_id": task_id}


def _run_extract(task_id: str, tmp: Path, kind: str) -> None:
    """后台识别线程：markitdown → OCR 兜底 → LLM 总结。各阶段检查取消标记。"""
    t = EXTRACT_TASKS[task_id]
    suffix = tmp.suffix.lower()
    try:
        text, ocr_used = _doc_to_text(tmp, suffix, cancel_check=lambda: t.get("cancel", False))
        t["ocr"] = ocr_used
        if t.get("cancel"):
            raise _ExtractCancelled()
        if kind == "template" and text.strip():
            # 话术是用户原始资产：只做文档解析/OCR，不交给 LLM 改写。
            items = [{
                "title": Path(t["filename"] or "话术模板").stem[:80],
                "content": text.strip(),
                "tags": ["邮件", "模板"],
            }]
            mode = "raw"
        else:
            items = _parse_product_doc(text)
            mode = "raw"
        if items and kind != "template":
            summary = _summarize_doc(
                text, Path(t["filename"] or "文档").stem[:40], kind,
                cancel_check=lambda: t.get("cancel", False),
            )
            if summary:
                items = [summary]
                mode = "summary"
        if not items:
            t["status"] = "error"
            t["error"] = "未识别到内容，文档可能没有可读文字"
            t["finished_at"] = time.time()
            return
        t["items"] = items
        t["mode"] = mode
        t["status"] = "done"
        t["finished_at"] = time.time()
    except _ExtractCancelled:
        t["status"] = "cancelled"
        t["finished_at"] = time.time()
    except ImportError:
        t["status"] = "error"
        t["error"] = "缺少识别库，请先运行: pip install markitdown python-multipart rapidocr-onnxruntime pymupdf"
        t["finished_at"] = time.time()
    except Exception as exc:  # noqa: BLE001
        import traceback

        err_log = DATA_DIR / "ocr_errors.log"
        try:
            with open(err_log, "a", encoding="utf-8") as f:
                f.write(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] {kind} {t['filename']}\n{traceback.format_exc()}\n")
        except OSError:
            pass
        t["status"] = "error"
        t["error"] = f"识别失败[{type(exc).__name__}]: {exc}"
        t["finished_at"] = time.time()
    finally:
        tmp.unlink(missing_ok=True)


@app.post("/api/products/extract")
async def api_product_extract(file: UploadFile = File(...)) -> dict:
    """文档识别导入（产品）：PDF/Word/Excel/PPT/Markdown/HTML → 产品条目。

    ① markitdown 转文本；提取不到文字（图片型 PDF）自动切 OCR；
    ② LLM 自动总结成一条结构化产品资料（有 key 时）；失败回退原始识别块。
    后台任务执行，前端轮询 GET /api/extract/tasks/{id}，可取消。
    """
    return await _start_extract(file, "product")


@app.post("/api/templates/extract")
async def api_template_extract(file: UploadFile = File(...)) -> dict:
    """历史邮件/聊天记录 → 话术模板条目（解析 + LLM 总结为风格样本）。"""
    return await _start_extract(file, "template")


@app.get("/api/extract/tasks/{task_id}")
def api_extract_task(task_id: str) -> dict:
    t = EXTRACT_TASKS.get(task_id)
    if not t:
        return {"status": "not_found"}
    return {k: t[k] for k in ("status", "error", "items", "mode", "ocr", "filename") if k in t}


@app.post("/api/extract/tasks/{task_id}/cancel")
def api_extract_cancel(task_id: str) -> dict:
    """请求停止识别：置取消标记，OCR 页循环/LLM 调用前检查即停（页粒度）。"""
    t = EXTRACT_TASKS.get(task_id)
    if not t:
        return {"ok": False, "error": "任务不存在"}
    if t["status"] != "running":
        return {"ok": False, "error": "任务已结束"}
    t["cancel"] = True
    return {"ok": True}


# ---------- 业务 API：照片库 + 照片筛选（任务制） ----------

PHOTO_SUFFIXES = {".jpg", ".jpeg", ".png", ".webp", ".gif"}
PHOTO_TASKS: dict[str, dict] = {}
PHOTO_TASK_LOCK = threading.Lock()
PHOTO_RUN_LOCK = threading.Lock()


class PhotoInput(BaseModel):
    filename: str = "photo.jpg"
    data_base64: str


class PhotoScanRequest(BaseModel):
    images: list[PhotoInput]


class PhotoLibraryItem(BaseModel):
    store_id: str
    photo_id: str


class PhotoDeleteRequest(BaseModel):
    items: list[PhotoLibraryItem]


class PhotoStoreRenameRequest(BaseModel):
    store_id: str
    name: str


class PhotoStoreDeleteRequest(BaseModel):
    store_id: str


def _photo_token(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()[:12]


def _find_library_store(store_id: str) -> Path | None:
    if not CRAWLER_PHOTOS_DIR.exists():
        return None
    root = CRAWLER_PHOTOS_DIR.resolve()
    for folder in CRAWLER_PHOTOS_DIR.iterdir():
        if not folder.is_dir() or _photo_token(folder.name) != store_id:
            continue
        candidate = folder.resolve()
        if root in candidate.parents:
            return candidate
    return None


def _photo_library_snapshot() -> list[dict]:
    """从爬虫照片目录增量生成只读索引；不创建或覆盖任何用户数据文件。"""
    if not CRAWLER_PHOTOS_DIR.exists():
        return []
    stores: list[dict] = []
    for folder in sorted((p for p in CRAWLER_PHOTOS_DIR.iterdir() if p.is_dir()), key=lambda p: p.name.casefold()):
        store_id = _photo_token(folder.name)
        photos = []
        for path in sorted(folder.iterdir(), key=lambda p: p.name.casefold()):
            if not path.is_file() or path.suffix.lower() not in PHOTO_SUFFIXES:
                continue
            stat = path.stat()
            photo_id = _photo_token(f"{folder.name}/{path.name}")
            photos.append({
                "photo_id": photo_id,
                "filename": path.name,
                "size": stat.st_size,
                "modified_at": stat.st_mtime,
                "url": f"/api/photos/library/image/{store_id}/{photo_id}",
            })
        if photos:
            stores.append({"store_id": store_id, "name": folder.name, "count": len(photos), "photos": photos})
    return stores


def _find_library_photo(store_id: str, photo_id: str) -> Path | None:
    folder = _find_library_store(store_id)
    if folder is None:
        return None
    root = CRAWLER_PHOTOS_DIR.resolve()
    for path in folder.iterdir():
        if not path.is_file() or path.suffix.lower() not in PHOTO_SUFFIXES:
            continue
        if _photo_token(f"{folder.name}/{path.name}") != photo_id:
            continue
        candidate = path.resolve()
        if root in candidate.parents:
            return candidate
    return None


@app.get("/api/photos/library")
def api_photo_library() -> dict:
    stores = _photo_library_snapshot()
    return {"stores": stores, "total": sum(store["count"] for store in stores)}


@app.get("/api/photos/library/image/{store_id}/{photo_id}")
def api_photo_library_image(store_id: str, photo_id: str):
    path = _find_library_photo(store_id, photo_id)
    if path is None:
        return {"error": "照片不存在"}
    return FileResponse(path)


@app.post("/api/photos/library/delete")
def api_photo_library_delete(req: PhotoDeleteRequest) -> dict:
    removed = 0
    affected: set[Path] = set()
    for item in req.items[:100]:
        path = _find_library_photo(item.store_id, item.photo_id)
        if path is None:
            continue
        try:
            affected.add(path.parent)
            path.unlink()
            removed += 1
        except OSError:
            continue
    for folder in affected:
        try:
            if folder.exists() and not any(folder.iterdir()):
                folder.rmdir()
        except OSError:
            pass
    return {"ok": True, "removed": removed}


@app.post("/api/photos/library/store/rename")
def api_photo_store_rename(req: PhotoStoreRenameRequest) -> dict:
    folder = _find_library_store(req.store_id)
    if folder is None:
        return {"ok": False, "error": "店铺不存在"}
    safe_name = re.sub(r'[\\/:*?"<>|\x00-\x1f]', "_", req.name).strip(" .")[:80]
    if not safe_name:
        return {"ok": False, "error": "店铺名称不能为空"}
    destination = folder.parent / safe_name
    same_folder = destination.exists() and destination.samefile(folder)
    if destination.exists() and not same_folder:
        return {"ok": False, "error": "已有同名店铺，请换一个名称"}
    try:
        if destination != folder:
            folder.rename(destination)
    except OSError as exc:
        return {"ok": False, "error": f"重命名失败：{exc}"}
    return {"ok": True, "store_id": _photo_token(safe_name), "name": safe_name}


@app.post("/api/photos/library/store/delete")
def api_photo_store_delete(req: PhotoStoreDeleteRequest) -> dict:
    folder = _find_library_store(req.store_id)
    if folder is None:
        return {"ok": False, "removed": 0, "error": "店铺不存在"}
    removed = 0
    try:
        for path in folder.iterdir():
            if path.is_file() and path.suffix.lower() in PHOTO_SUFFIXES:
                path.unlink()
                removed += 1
        if not any(folder.iterdir()):
            folder.rmdir()
    except OSError as exc:
        return {"ok": False, "removed": removed, "error": f"删除失败：{exc}"}
    if folder.exists():
        return {"ok": True, "removed": removed, "warning": "目录中的非照片文件已保留"}
    return {"ok": True, "removed": removed}


def _decode_photo(data_base64: str) -> bytes:
    encoded = data_base64.split(",", 1)[-1].strip()
    return base64.b64decode(encoded, validate=True)


def _photo_suffix(data: bytes) -> str:
    if data.startswith(b"\x89PNG"):
        return ".png"
    if data.startswith(b"RIFF") and data[8:12] == b"WEBP":
        return ".webp"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return ".gif"
    return ".jpg"


def _safe_photo_filename(filename: str, data: bytes) -> str:
    stem = Path(filename or "photo").stem
    stem = re.sub(r'[\\/:*?"<>|\x00-\x1f]', "_", stem).strip(" .")[:60] or "photo"
    return f"{stem}{_photo_suffix(data)}"


def _prepare_photo_for_vision(data: bytes) -> bytes:
    """长边压到 1024px，减少视觉 token；Pillow 异常时保留原图。"""
    try:
        from PIL import Image, ImageOps

        with Image.open(io.BytesIO(data)) as image:
            image = ImageOps.exif_transpose(image)
            image.thumbnail((1024, 1024))
            if image.mode not in ("RGB", "L"):
                background = Image.new("RGB", image.size, "white")
                if "A" in image.getbands():
                    background.paste(image, mask=image.getchannel("A"))
                else:
                    background.paste(image.convert("RGB"))
                image = background
            elif image.mode == "L":
                image = image.convert("RGB")
            output = io.BytesIO()
            image.save(output, format="JPEG", quality=84, optimize=True)
            return output.getvalue()
    except Exception:  # noqa: BLE001 - 压缩失败不阻塞识别
        return data


def _finish_photo_task(task: dict, status: str, error: str = "") -> None:
    task["status"] = status
    task["error"] = error
    task["finished_at"] = time.time()


def _run_photo_scan(task_id: str, images: list[PhotoInput]) -> None:
    from app.crawler.vision_analyzer import analyze_image_bytes_with_meta

    task = PHOTO_TASKS[task_id]
    PHOTO_SCAN_DIR.mkdir(parents=True, exist_ok=True)
    try:
        for index, image in enumerate(images, start=1):
            if task.get("cancel"):
                _finish_photo_task(task, "cancelled")
                return
            result = {
                "filename": image.filename,
                "saved_path": "",
                "has_device": False,
                "confidence": 0.0,
                "devices": [],
                "error": None,
            }
            try:
                data = _decode_photo(image.data_base64)
                if not data or len(data) > 10 * 1024 * 1024:
                    raise ValueError("图片为空或超过 10MB")
                safe_name = _safe_photo_filename(image.filename, data)
                saved = PHOTO_SCAN_DIR / f"{task_id[:8]}_{index:02d}_{safe_name}"
                saved.write_bytes(data)
                result["saved_path"] = str(saved.relative_to(ROOT))
                prepared = _prepare_photo_for_vision(data)
                items, error, used = analyze_image_bytes_with_meta(
                    [prepared], proxy=settings.crawler_proxy, verbose=True,
                )
                result["devices"] = items
                result["has_device"] = bool(items)
                result["confidence"] = max((float(item.get("confidence", 0)) for item in items), default=0.0)
                result["error"] = error
                if used:
                    result["provider"] = used["provider"]
                    result["model"] = used["model"]
            except Exception as exc:  # noqa: BLE001 - 单张失败继续下一张
                result["error"] = str(exc)
            with PHOTO_TASK_LOCK:
                result["preview_url"] = f"/api/photos/tasks/{task_id}/image/{len(task['results'])}"
                task["results"].append(result)
                task["done"] = index
        _finish_photo_task(task, "done")
    except Exception as exc:  # noqa: BLE001
        _finish_photo_task(task, "error", str(exc))
    finally:
        PHOTO_RUN_LOCK.release()


@app.post("/api/photos/start")
def api_photos_start(req: PhotoScanRequest) -> dict:
    if not req.images:
        return {"task_id": "", "error": "请先选择照片"}
    if len(req.images) > 50:
        return {"task_id": "", "error": "单次最多选择 50 张照片"}
    estimated_bytes = sum(len(image.data_base64.split(",", 1)[-1]) * 3 // 4 for image in req.images)
    if estimated_bytes > 50 * 1024 * 1024:
        return {"task_id": "", "error": "单次照片总大小不能超过 50MB"}
    current_vision = get_vision_config()
    has_vision_candidate = (
        bool(get_vision_failover_configs(current_vision["provider"]))
        if current_vision.get("provider") else bool(current_vision.get("api_key"))
    )
    if not has_vision_candidate:
        return {"task_id": "", "error": "请先在设置中配置视觉识别 API Key"}
    if not PHOTO_RUN_LOCK.acquire(blocking=False):
        return {"task_id": "", "error": "已有照片识别任务正在运行，请等待完成或先取消"}
    task_id = uuid.uuid4().hex[:12]
    with PHOTO_TASK_LOCK:
        prune_finished_tasks(PHOTO_TASKS)
        PHOTO_TASKS[task_id] = {
            "status": "running", "cancel": False, "total": len(req.images), "done": 0,
            "results": [], "error": "", "created_at": time.time(),
        }
    threading.Thread(target=_run_photo_scan, args=(task_id, req.images), daemon=True).start()
    return {"task_id": task_id}


@app.get("/api/photos/tasks/{task_id}")
def api_photos_task(task_id: str) -> dict:
    task = PHOTO_TASKS.get(task_id)
    if not task:
        return {"status": "not_found", "total": 0, "done": 0, "results": []}
    return {key: task.get(key) for key in ("status", "total", "done", "results", "error")}


@app.get("/api/photos/tasks/{task_id}/image/{result_index}")
def api_photos_task_image(task_id: str, result_index: int):
    task = PHOTO_TASKS.get(task_id)
    results = task.get("results", []) if task else []
    if result_index < 0 or result_index >= len(results):
        return {"error": "照片不存在"}
    saved_path = results[result_index].get("saved_path", "")
    if not saved_path:
        return {"error": "照片不存在"}
    path = (ROOT / saved_path).resolve()
    root = PHOTO_SCAN_DIR.resolve()
    if root not in path.parents or not path.is_file():
        return {"error": "照片不存在"}
    return FileResponse(path)


@app.post("/api/photos/tasks/{task_id}/cancel")
def api_photos_cancel(task_id: str) -> dict:
    task = PHOTO_TASKS.get(task_id)
    if not task:
        return {"ok": False, "error": "任务不存在"}
    if task["status"] != "running":
        return {"ok": False, "error": "任务已结束"}
    task["cancel"] = True
    return {"ok": True}


@app.put("/api/products/{item_id}")
def api_product_update(item_id: str, req: ProductRequest) -> dict:
    items = _read_json(PRODUCTS_FILE, [])
    for it in items:
        if it.get("id") == item_id:
            it["title"] = req.title.strip()
            it["content"] = req.content.strip()
            if req.tags:
                it["tags"] = req.tags
            _write_json(PRODUCTS_FILE, items)
            return {"ok": True, "item": it}
    return {"ok": False, "error": "产品不存在"}


@app.delete("/api/products/{item_id}")
def api_product_delete(item_id: str) -> dict:
    items = _read_json(PRODUCTS_FILE, [])
    items = [it for it in items if it.get("id") != item_id]
    _write_json(PRODUCTS_FILE, items)
    return {"ok": True}


# ---------- 业务 API：邮件模板 ----------

@app.get("/api/templates")
def api_templates() -> list[dict]:
    return _read_json(TEMPLATES_FILE, [])


@app.post("/api/templates")
def api_template_add(req: ProductRequest) -> dict:
    items = _read_json(TEMPLATES_FILE, [])
    title = req.title.strip()
    content = req.content.strip()
    if not title or not content:
        return {"ok": False, "error": "标题和内容必填"}
    normalized = re.sub(r"\s+", " ", content).strip().casefold()
    for existing in items:
        existing_content = re.sub(r"\s+", " ", str(existing.get("content") or "")).strip().casefold()
        if existing_content == normalized:
            return {"ok": True, "item": existing, "duplicate": True}
    item = {
        "id": _next_id(items, "eml-"),
        "title": title,
        "content": content,
        "source": "frontend",
        "tags": req.tags or ["邮件", "模板"],
    }
    items.append(item)
    _write_json(TEMPLATES_FILE, items)
    return {"ok": True, "item": item}


@app.delete("/api/templates/{item_id}")
def api_template_delete(item_id: str) -> dict:
    items = _read_json(TEMPLATES_FILE, [])
    items = [it for it in items if it.get("id") != item_id]
    _write_json(TEMPLATES_FILE, items)
    return {"ok": True}


@app.put("/api/templates/{item_id}")
def api_template_update(item_id: str, req: ProductRequest) -> dict:
    items = _read_json(TEMPLATES_FILE, [])
    for it in items:
        if it.get("id") == item_id:
            it["title"] = req.title.strip()
            it["content"] = req.content.strip()
            if req.tags:
                it["tags"] = req.tags
            _write_json(TEMPLATES_FILE, items)
            return {"ok": True, "item": it}
    return {"ok": False, "error": "模板不存在"}


# ---------- 业务 API：客户名单 ----------

def _customer_identity_keys(customer: dict) -> set[str]:
    """生成稳定的客户查重键；任一可靠联系方式或“名称+地区”相同即视为重复。"""
    keys: set[str] = set()

    email = str(customer.get("email") or "").strip().casefold()
    if "@" in email:
        keys.add(f"email:{email}")

    phone = re.sub(r"\D", "", str(customer.get("phone") or ""))
    if len(phone) >= 7:
        keys.add(f"phone:{phone}")

    website = str(customer.get("website") or "").strip().casefold()
    if website:
        parsed = urlsplit(website if "://" in website else f"https://{website}")
        host = (parsed.hostname or "").removeprefix("www.")
        if host and host not in {"wa.me", "api.whatsapp.com"}:
            keys.add(f"website:{host}")

    def compact(value: object) -> str:
        return re.sub(r"[^\w]+", "", str(value or "").casefold())

    name = compact(customer.get("name"))
    if name:
        country = compact(customer.get("country"))
        city = compact(customer.get("city"))
        keys.add(f"name:{name}|{country}|{city}")
    return keys


def _customer_seen_keys(items: list[dict]) -> set[str]:
    return {key for item in items for key in _customer_identity_keys(item)}

@app.get("/api/customers")
def api_customers() -> list[dict]:
    return _read_json(CUSTOMERS_FILE, [])


@app.post("/api/customers")
def api_customer_add(req: CustomerRequest) -> dict:
    items = _read_json(CUSTOMERS_FILE, [])
    item = {
        "id": _next_id(items, "c-"),
        "name": req.name.strip(),
        "country": req.country.strip(),
        "city": req.city.strip(),
        "email": req.email.strip(),
        "website": req.website.strip(),
        "notes": req.notes.strip(),
        "source": "manual",
        "status": req.status if req.status in ("new", "sent", "replied") else "new",
        "created_at": time.time(),
    }
    if _customer_identity_keys(item) & _customer_seen_keys(items):
        return {"ok": False, "duplicate": True, "error": "该客户已在名单中"}
    items.append(item)
    _write_json(CUSTOMERS_FILE, items)
    return {"ok": True, "item": item}


@app.put("/api/customers/{item_id}")
def api_customer_update(item_id: str, req: CustomerRequest) -> dict:
    items = _read_json(CUSTOMERS_FILE, [])
    for it in items:
        if it.get("id") == item_id:
            it.update({
                "name": req.name.strip(), "country": req.country.strip(),
                "city": req.city.strip(), "email": req.email.strip(),
                "website": req.website.strip(), "phone": req.phone.strip(),
                "notes": req.notes.strip(),
                "status": req.status if req.status in ("new", "sent", "replied") else it.get("status", "new"),
            })
            _write_json(CUSTOMERS_FILE, items)
            return {"ok": True, "item": it}
    return {"ok": False, "error": "客户不存在"}


@app.delete("/api/customers/{item_id}")
def api_customer_delete(item_id: str) -> dict:
    items = _read_json(CUSTOMERS_FILE, [])
    items = [it for it in items if it.get("id") != item_id]
    _write_json(CUSTOMERS_FILE, items)
    return {"ok": True}


@app.post("/api/customers/batch-delete")
def api_customer_batch_delete(body: dict) -> dict:
    ids = {str(item_id) for item_id in (body.get("ids") or []) if item_id}
    if not ids:
        return {"ok": False, "error": "未选择客户"}
    items = _read_json(CUSTOMERS_FILE, [])
    kept = [item for item in items if str(item.get("id")) not in ids]
    removed = len(items) - len(kept)
    if removed:
        _write_json(CUSTOMERS_FILE, kept)
    return {"ok": True, "removed": removed, "total": len(kept)}


@app.post("/api/customers/import")
def api_customer_import(body: dict) -> dict:
    """导入 CSV 文本（utf-8-sig 兼容 Excel），字段 name,country,city,email,website,notes。
    source 标记来源：manual（手动/CSV）| maps | webs，客户名单按此分组。"""
    text = (body.get("text") or "").strip()
    src = (body.get("source") or "manual").strip()
    if src not in ("manual", "maps", "webs"):
        src = "manual"
    if not text:
        return {"ok": False, "error": "内容为空"}
    try:
        rows = list(csv.DictReader(io.StringIO(text)))
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"CSV 解析失败: {exc}"}

    items = _read_json(CUSTOMERS_FILE, [])
    seen = _customer_seen_keys(items)
    added = 0
    duplicates = 0
    invalid = 0
    added_items: list[dict] = []
    next_number = int(_next_id(items, "c-").removeprefix("c-"))
    for row in rows:
        name = (row.get("name") or "").strip()
        if not name:
            invalid += 1
            continue
        email_raw = (row.get("email") or "").strip()
        # 爬虫占位符（（WhatsApp 跟进）等）不算真邮箱，置空避免脏数据
        if "WhatsApp" in email_raw or "跟进" in email_raw:
            email_raw = ""
        item = {
            "id": f"c-{next_number:03d}",
            "name": name,
            "country": (row.get("country") or "").strip(),
            "city": (row.get("city") or "").strip(),
            "email": email_raw,
            "website": (row.get("website") or "").strip(),
            "phone": (row.get("phone") or "").strip(),
            # 爬虫 CSV 的 instrument/instruments 列并入 notes，批量生成时自动作为客户背景
            "notes": _merge_notes(row),
            "source": src,
            "status": "new",
            "created_at": time.time(),
        }
        keys = _customer_identity_keys(item)
        if keys & seen:
            duplicates += 1
            continue
        seen.update(keys)
        items.append(item)
        added_items.append(item)
        added += 1
        next_number += 1
    # 全部重复时不再重写整份名单；这是重复点击导入时最常见的慢路径。
    if added:
        _write_json(CUSTOMERS_FILE, items)
    return {
        "ok": True,
        "added": added,
        "duplicates": duplicates,
        "invalid": invalid,
        "total": len(items),
        "items": added_items,
    }


def _merge_notes(row: dict) -> str:
    """客户导入时合并备注、仪器和缺品推荐，并对列表内容归一化去重。"""
    parts: list[str] = []
    notes = (row.get("notes") or "").strip()
    if notes:
        parts.append(notes)

    def normalized_list(value: object, separators: str, joiner: str) -> str:
        seen: set[str] = set()
        values: list[str] = []
        for item in re.split(separators, str(value or "")):
            clean = item.strip()
            key = clean.casefold()
            if clean and key not in seen:
                seen.add(key)
                values.append(clean)
        return joiner.join(values)

    ins = (row.get("instruments") or "").strip()
    if not ins:
        ins = (row.get("instrument") or "").strip()
    if ins.casefold() not in ("", "yes", "no", "unknown", "massage"):
        clean_instruments = normalized_list(ins, r"[|,，;；]+", "、")
        detail = f"店舖用儀器：{clean_instruments}"
        if clean_instruments and detail not in notes:
            parts.append(detail)

    gap_recs = normalized_list(row.get("gap_recs"), r"[|;；]+", "；")
    if gap_recs:
        detail = f"缺品推薦：{gap_recs}"
        if detail not in notes:
            parts.append(detail)
    return "；".join(parts)


# ---------- 业务 API：历史记录 ----------

@app.get("/api/history")
def api_history(limit: int = 30) -> list[dict]:
    return recent_emails(limit=min(limit, 200))


@app.delete("/api/history/{item_id}")
def api_history_delete(item_id: int) -> dict:
    delete_email(item_id)
    return {"ok": True}


@app.post("/api/history/batch-delete")
def api_history_batch_delete(body: dict) -> dict:
    ids: list[int] = []
    for raw_id in body.get("ids") or []:
        try:
            ids.append(int(raw_id))
        except (TypeError, ValueError):
            continue
    if not ids:
        return {"ok": False, "error": "未选择历史记录"}
    return {"ok": True, "removed": delete_emails(ids)}


# ---------- 业务 API：爬虫 ----------

CRAWLER_TASKS: dict[str, dict] = {}
CRAWLER_RUN_LOCK = threading.Lock()


def _append_crawler_log(task: dict, line: str) -> None:
    """追加任务进度，避免 redirect_stdout 修改整个进程的标准输出。"""
    if not line.strip():
        return
    log = task.setdefault("log", [])
    log.append(line)
    task["log"] = log[-300:]


class CrawlerRequest(BaseModel):
    queries: str = "medspa,medical spa,醫學美容"
    country: str = "香港"
    targets: list[str] = ["all"]  # all | email | phone | whatsapp（可组合）
    max_customers: int = 20
    source: str = "maps"  # maps（Google Maps 详情） | webs（官网深挖）
    engine: str = "auto"  # webs 用：auto | google | bing | ddg
    region: str = "hk"    # webs 用：hk（只收 .hk 域名） | any
    allow_recrawl: bool = False  # Maps/Webs：允许返回已经爬过或导入过的客户
    save_photos: bool = True  # Maps：视觉识别后保存照片，供人工抽查


def _website_domain(value: object) -> str:
    raw = str(value or "").strip().casefold()
    if not raw:
        return ""
    parsed = urlsplit(raw if "://" in raw else f"https://{raw}")
    return (parsed.hostname or "").removeprefix("www.")


def _webs_excluded_domains() -> set[str]:
    """默认排除客户名单和历次 Webs 成功结果，避免重复进入同一官网。"""
    domains = {
        domain for item in _read_json(CUSTOMERS_FILE, [])
        if (domain := _website_domain(item.get("website")))
    }
    for item in _read_json(CRAWLER_SEEN_FILE, []):
        domain = _website_domain(item.get("domain") if isinstance(item, dict) else item)
        if domain:
            domains.add(domain)
    return domains


def _crawler_match_keys(item: dict) -> set[str]:
    """爬虫跨来源查重键：官网、电话，以及标准化名称+国家。"""
    keys: set[str] = set()
    domain = _website_domain(item.get("website") or item.get("domain"))
    if domain and domain not in {"wa.me", "api.whatsapp.com", "bit.ly"}:
        keys.add(f"website:{domain}")
    phone = re.sub(r"\D", "", str(item.get("phone") or ""))
    if len(phone) >= 7:
        keys.add(f"phone:{phone}")

    def compact(value: object) -> str:
        return re.sub(r"[^\w]+", "", str(value or "").casefold())

    name = compact(item.get("name"))
    if name:
        keys.add(f"name:{name}|{compact(item.get('country'))}")
    return keys


def _maps_exclusions() -> tuple[set[str], set[str]]:
    """读取 Maps 地点键，并合并客户名单/爬取历史的身份键。"""
    history = _read_json(CRAWLER_SEEN_FILE, [])
    place_keys = {
        str(item.get("place_key")) for item in history
        if isinstance(item, dict) and item.get("place_key")
    }
    identities = {
        key for item in [*_read_json(CUSTOMERS_FILE, []), *history]
        if isinstance(item, dict)
        for key in _crawler_match_keys(item)
    }
    return place_keys, identities


def _remember_webs_results(results: list[dict]) -> None:
    """记录已成功返回的官网；用户即使不导入，下次也不会再次拿到。"""
    items = _read_json(CRAWLER_SEEN_FILE, [])
    known = {
        domain for item in items
        if (domain := _website_domain(item.get("domain") if isinstance(item, dict) else item))
    }
    now = time.time()
    changed = False
    for result in results:
        domain = _website_domain(result.get("website"))
        if domain and domain not in known:
            items.append({"domain": domain, "first_seen_at": now})
            known.add(domain)
            changed = True
    if changed:
        _write_json(CRAWLER_SEEN_FILE, items[-5000:])


def _remember_maps_results(results: list[dict]) -> None:
    """持久化 Maps 成功结果的地点键和辅助身份，供下次在深挖前排除。"""
    items = _read_json(CRAWLER_SEEN_FILE, [])
    known_places = {
        str(item.get("place_key")) for item in items
        if isinstance(item, dict) and item.get("place_key")
    }
    now = time.time()
    changed = False
    for result in results:
        place_key = str(result.get("_place_key") or "")
        if not place_key or place_key in known_places:
            continue
        items.append({
            "source": "maps", "place_key": place_key,
            "name": str(result.get("name") or ""),
            "country": str(result.get("country") or ""),
            "domain": _website_domain(result.get("website")),
            "website": str(result.get("website") or ""),
            "phone": str(result.get("phone") or ""),
            "first_seen_at": now,
        })
        known_places.add(place_key)
        changed = True
    if changed:
        _write_json(CRAWLER_SEEN_FILE, items[-5000:])


def _filter_crawler_results(results: list[dict], targets: list[str]) -> list[dict]:
    """按目标类型过滤，并保持原有结果顺序。"""
    return [result for result in results if matches_targets(result, targets)]


@app.post("/api/crawler/start")
def api_crawler_start(req: CrawlerRequest) -> dict:
    """创建爬虫任务：Google Maps 搜索多关键词 → 官网/邮箱/电话/WhatsApp → 结果回传。"""
    if not CRAWLER_RUN_LOCK.acquire(blocking=False):
        return {"task_id": "", "error": "已有爬虫任务正在运行，请等待完成或先取消"}
    queries = [q.strip() for q in req.queries.split(",") if q.strip()] or ["beauty salon"]
    task_id = uuid.uuid4().hex[:12]
    prune_finished_tasks(CRAWLER_TASKS)
    CRAWLER_TASKS[task_id] = {
        "status": "running", "cancel": False, "log": [], "results": [], "error": "",
        "created_at": time.time(),
    }

    def execute() -> None:
        import traceback as _tb

        task = CRAWLER_TASKS[task_id]
        try:
            with use_progress_sink(lambda line: _append_crawler_log(task, line)):
                if req.source == "webs":
                    # Webs：网页搜索 → 官网深挖邮箱/电话/WhatsApp/社媒/仪器品牌
                    from app.crawler.webs_hunter import hunt_websites
                    excluded_domains = set() if req.allow_recrawl else _webs_excluded_domains()
                    if excluded_domains:
                        _append_crawler_log(task, f"[webs] 已加载 {len(excluded_domains)} 个历史域名，本次自动跳过")
                    raw = hunt_websites(
                        queries, max_customers=req.max_customers, engine=req.engine,
                        region=req.region, verbose=True,
                        cancel_check=lambda: CRAWLER_TASKS[task_id].get("cancel", False),
                        result_filter=lambda item: matches_targets(item, req.targets),
                        excluded_domains=excluded_domains,
                    )
                    _remember_webs_results(raw)
                    results = []
                    for r in raw:
                        wa = r.get("whatsapp", "") or ""
                        digits = re.sub(r"\D", "", wa)
                        results.append({
                            "name": r.get("name", ""), "country": req.country, "city": req.country,
                            "website": r.get("website", ""), "email": r.get("email", ""),
                            "phone": r.get("phone", ""),
                            "wa_link": f"https://wa.me/{digits}" if digits else "",
                            "whatsapp": wa,
                            "facebook": r.get("facebook", ""),
                            "instagram": r.get("instagram", ""),
                            "instrument": "yes" if r.get("instruments") else "",
                            "instruments": r.get("instruments", ""),
                            "gap_recs": r.get("gap_recs", ""),
                        })
                else:
                    # Maps：Google Maps 详情页（原逻辑）
                    from app.crawler.maps_hunter import hunt_maps_customers
                    _append_crawler_log(
                        task,
                        f"[maps] 视觉照片落盘：{'已开启' if req.save_photos else '已关闭'}",
                    )
                    excluded_places, excluded_identities = (set(), set()) if req.allow_recrawl else _maps_exclusions()
                    if excluded_places or excluded_identities:
                        _append_crawler_log(
                            task,
                            f"[maps] 已加载 {len(excluded_places)} 个历史地点、{len(excluded_identities)} 个客户标识，本次自动跳过",
                        )
                    raw = hunt_maps_customers(
                        queries, req.country, max_customers=req.max_customers,
                        headless=True, verbose=True,
                        cancel_check=lambda: CRAWLER_TASKS[task_id].get("cancel", False),
                        result_filter=lambda item: matches_targets(item, req.targets),
                        excluded_place_keys=excluded_places,
                        exclude_filter=(
                            (lambda item: bool(_crawler_match_keys(item) & excluded_identities))
                            if excluded_identities else None
                        ),
                        save_photos=req.save_photos,
                    )
                    _remember_maps_results(raw)
                    results = [
                        {key: value for key, value in item.items() if key != "_place_key"}
                        for item in raw
                    ]
        except Exception as exc:  # noqa: BLE001
            task["error"] = str(exc)
            task["status"] = "error"
            task["finished_at"] = time.time()
            # 完整 traceback 落盘，网页 toast 只显示一行，诊断靠日志
            try:
                with open(DATA_DIR / "crawler_errors.log", "a", encoding="utf-8") as f:
                    f.write(f"\n--- {time.strftime('%Y-%m-%d %H:%M:%S')} {task_id} ---\n")
                    f.write(_tb.format_exc())
            except Exception:  # noqa: BLE001
                pass
            return
        if task.get("cancel"):
            task["status"] = "cancelled"
            task["finished_at"] = time.time()
            return
        task["results"] = _filter_crawler_results(results, req.targets)
        task["status"] = "done"
        task["finished_at"] = time.time()

    def run() -> None:
        try:
            execute()
        finally:
            CRAWLER_RUN_LOCK.release()

    try:
        threading.Thread(target=run, daemon=True).start()
    except Exception:
        CRAWLER_RUN_LOCK.release()
        CRAWLER_TASKS.pop(task_id, None)
        raise
    return {"task_id": task_id}


@app.get("/api/crawler/tasks/{task_id}")
def api_crawler_task(task_id: str) -> dict:
    t = CRAWLER_TASKS.get(task_id)
    if not t:
        return {"status": "not_found"}
    return t


@app.post("/api/crawler/tasks/{task_id}/cancel")
def api_crawler_cancel(task_id: str) -> dict:
    """请求停止爬虫：置取消标记，当前店处理完即停（每家之间检查）。"""
    t = CRAWLER_TASKS.get(task_id)
    if not t:
        return {"ok": False, "error": "任务不存在"}
    if t["status"] != "running":
        return {"ok": False, "error": "任务已结束"}
    t["cancel"] = True
    return {"ok": True}


# ---------- 静态托管（API 路由定义之后再 mount，避免覆盖） ----------

if FRONTEND_DIST is not None and FRONTEND_DIST.exists():
    app.mount("/", StaticFiles(directory=str(FRONTEND_DIST), html=True), name="frontend")
